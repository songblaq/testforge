"""Office document parsing -- PowerPoint, Word, and Excel."""

from __future__ import annotations

from pathlib import Path
from typing import Any


def parse_file(path: Path) -> dict[str, Any]:
    """Parse an Office document and return structured content."""
    suffix = path.suffix.lower()

    if suffix in (".pptx", ".ppt"):
        return _parse_pptx(path)
    elif suffix in (".docx", ".doc"):
        return _parse_docx(path)
    elif suffix in (".xlsx", ".xls"):
        return _parse_xlsx(path)
    else:
        raise ValueError(f"Unsupported office format: {suffix}")


def _parse_pptx(path: Path) -> dict[str, Any]:
    """Parse a PowerPoint file."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slides.append({"slide": idx, "texts": texts})

    return {
        "type": "pptx",
        "source": str(path),
        "slides": slides,
        "slide_count": len(slides),
    }


def _parse_docx(path: Path) -> dict[str, Any]:
    """Parse a Word document."""
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    return {
        "type": "docx",
        "source": str(path),
        "paragraphs": paragraphs,
        "paragraph_count": len(paragraphs),
    }


def _parse_xlsx(path: Path) -> dict[str, Any]:
    """Parse an Excel workbook."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets: list[dict[str, Any]] = []

    for name in wb.sheetnames:
        ws = wb[name]
        rows: list[list[Any]] = []
        for row in ws.iter_rows(values_only=True):
            rows.append(list(row))
        sheets.append({"name": name, "rows": rows, "row_count": len(rows)})

    wb.close()
    return {
        "type": "xlsx",
        "source": str(path),
        "sheets": sheets,
        "sheet_count": len(sheets),
    }
